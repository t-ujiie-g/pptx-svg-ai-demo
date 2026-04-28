"""Subprocess wrapper for pptx-svg's inspect script (PNG rendering).

backend/skills/pptx/scripts/pptx_inspect.js renders each slide of a PPTX to
PNG and extracts the shape structure, so the chat request can attach a visual
+ structural context summary to the user message.

Editing itself lives in the skill as Python (scripts/edit_pptx.py, invoked via
ADK's run_skill_script), so no Node-side edit path is needed here.
"""

from __future__ import annotations

import asyncio
import base64
import io
import json
import logging
import os
import pathlib
from typing import Any

from pptx import Presentation
from pptx.oxml.ns import qn

from src.constants import PPTX_INSPECT_TIMEOUT

logger = logging.getLogger(__name__)

_SCRIPTS_DIR = (
    pathlib.Path(__file__).resolve().parent.parent.parent
    / "skills" / "pptx" / "scripts"
)
_INSPECT_SCRIPT = _SCRIPTS_DIR / "pptx_inspect.js"
_NODE_PATH = os.environ.get("NODE_PATH", "/usr/lib/node_modules")


async def _run_node_script(
    script: pathlib.Path,
    stdin_data: bytes,
    extra_args: list[str] | None = None,
) -> dict[str, Any]:
    """Run a Node.js script, feed stdin, return parsed JSON from stdout."""
    cmd = ["node", str(script), *(extra_args or [])]
    env = {**os.environ, "NODE_PATH": _NODE_PATH}

    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdin=asyncio.subprocess.PIPE,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
        env=env,
    )
    try:
        stdout, stderr = await asyncio.wait_for(
            proc.communicate(input=stdin_data),
            timeout=PPTX_INSPECT_TIMEOUT,
        )
    except asyncio.TimeoutError:
        proc.kill()
        raise RuntimeError(f"{script.name} timed out after {PPTX_INSPECT_TIMEOUT}s")

    if proc.returncode != 0:
        err = stderr.decode("utf-8", errors="replace")[:2000]
        logger.error(f"{script.name} exit={proc.returncode}\nstderr:\n{err}")
        raise RuntimeError(f"{script.name} failed (exit {proc.returncode}): {err}")

    try:
        return json.loads(stdout)
    except json.JSONDecodeError as e:
        err = stderr.decode("utf-8", errors="replace")[:2000]
        head = stdout[:500].decode("utf-8", errors="replace")
        logger.error(
            f"{script.name} non-JSON stdout (len={len(stdout)}): "
            f"head={head!r}\nstderr:\n{err}"
        )
        raise RuntimeError(f"{script.name} produced non-JSON output: {e}")


def _cell_fill_hex(cell) -> str | None:
    """Return cell solid-fill colour as 6-hex (no `#`), or None if not a solid sRGB fill."""
    cell_pr = cell._tc.find(qn("a:tcPr"))
    if cell_pr is None:
        return None
    solid = cell_pr.find(qn("a:solidFill"))
    if solid is None:
        return None
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return None
    val = srgb.get("val")
    return val.upper() if val and len(val) == 6 else None


def _extract_tables(pptx_bytes: bytes) -> dict[int, dict[int, dict[str, Any]]]:
    """Walk the PPTX with python-pptx and pull out every table's structure.

    Returned mapping is keyed by (slide_idx, shape_idx). Shape index order matches
    pptx_inspect.js because both iterate <p:spTree> children in document order.
    """
    out: dict[int, dict[int, dict[str, Any]]] = {}
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as e:
        logger.warning(f"table extract: cannot open PPTX: {e}")
        return out

    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            cells: list[list[dict[str, Any]]] = []
            for row in table.rows:
                row_cells: list[dict[str, Any]] = []
                for cell in row.cells:
                    info: dict[str, Any] = {"text": cell.text}
                    fill_hex = _cell_fill_hex(cell)
                    if fill_hex:
                        info["fill_hex"] = fill_hex
                    row_cells.append(info)
                cells.append(row_cells)
            out.setdefault(s_idx, {})[sh_idx] = {
                "rows": len(table.rows),
                "cols": len(table.columns),
                "cells": cells,
            }
    return out


def _merge_tables(info: dict[str, Any], tables: dict[int, dict[int, dict[str, Any]]]) -> None:
    """Mutate `info` in place: attach extracted table data onto matching shape dicts."""
    for slide in info.get("slides", []):
        s_idx = slide.get("slide_idx")
        slide_tables = tables.get(s_idx)
        if not slide_tables:
            continue
        for shape in slide.get("shapes", []):
            tbl = slide_tables.get(shape.get("idx"))
            if tbl:
                shape["table"] = tbl


async def inspect_pptx(
    pptx_bytes: bytes,
    *,
    with_png: bool = True,
    png_width: int = 960,
    png_height: int = 540,
) -> dict[str, Any]:
    """Run pptx_inspect.js on the given PPTX bytes.

    Returns:
        {
            "slide_count": int,
            "slide_width_emu": int, "slide_height_emu": int,
            "slides": [
                {"slide_idx": int, "shapes": [...], "png_base64"?: str}
            ]
        }

    Each shape that backs a table gains a `"table": {rows, cols, cells: [[{text,
    fill_hex?}]]}` field, populated from a parallel python-pptx pass — the JS
    inspect renders SVG and so loses table-cell structure.
    """
    args = [f"--width={png_width}", f"--height={png_height}"]
    if not with_png:
        args.append("--no-png")
    b64 = base64.b64encode(pptx_bytes)
    info = await _run_node_script(_INSPECT_SCRIPT, b64, args)
    _merge_tables(info, _extract_tables(pptx_bytes))
    return info
